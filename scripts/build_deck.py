"""
Build the TwinForge business case deck (python-pptx).
Dark, product-matching (IDE/clay) theme. Screenshot placeholders left for the
user to fill. Run:  python -m scripts.build_deck
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

# ---- palette (matches the TwinForge UI) ----
BG    = RGBColor(0x14, 0x14, 0x16)
BG2   = RGBColor(0x1C, 0x1C, 0x20)
PANEL = RGBColor(0x22, 0x22, 0x28)
LINE  = RGBColor(0x3A, 0x3A, 0x42)
INK   = RGBColor(0xED, 0xED, 0xED)
INK2  = RGBColor(0x9E, 0x9E, 0xA6)
INK3  = RGBColor(0x6B, 0x6B, 0x74)
CLAY  = RGBColor(0xD9, 0x77, 0x57)
GREEN = RGBColor(0x5F, 0xB0, 0x6A)
AMBER = RGBColor(0xD9, 0xA4, 0x41)
RED   = RGBColor(0xD0, 0x68, 0x5C)
BLUE  = RGBColor(0x7F, 0xA8, 0xD8)

HEAD = "Calibri"
BODY = "Calibri"
MONO = "Consolas"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def slide(bg=BG):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = bg
    return s


def _set_dash(shape, dash="dash"):
    ln = shape.line._get_or_add_ln()
    d = ln.find(qn('a:prstDash'))
    if d is None:
        d = ln.makeelement(qn('a:prstDash'), {})
        ln.append(d)
    d.set('val', dash)


def box(s, x, y, w, h, fill=PANEL, line=LINE, line_w=0.75, radius=True):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    if radius:
        try:
            shp.adjustments[0] = 0.06
        except Exception:
            pass
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space=4, wrap=True):
    """runs: list of paragraphs; each paragraph is list of (txt, size, color, bold, font)."""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    for m in (tf.margin_left, ):
        pass
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space); p.space_before = Pt(0)
        for (txt, size, color, bold, font) in para:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.bold = bold
            r.font.color.rgb = color; r.font.name = font
    return tb


def R(txt, size=15, color=INK, bold=False, font=BODY):
    return (txt, size, color, bold, font)


def title(s, kicker, ttl):
    text(s, 0.6, 0.42, 12.1, 0.4, [[R(kicker.upper(), 12, CLAY, True, MONO)]])
    text(s, 0.6, 0.74, 12.1, 0.9, [[R(ttl, 32, INK, True, HEAD)]])


def bullets(s, x, y, w, h, items, size=14.5, gap=7, mark=CLAY, color=INK):
    runs = []
    for it in items:
        if isinstance(it, tuple):
            lead, rest = it
            runs.append([R("▸  ", size, mark, True, BODY), R(lead, size, INK, True, BODY),
                         R(rest, size, color, False, BODY)])
        else:
            runs.append([R("▸  ", size, mark, True, BODY), R(it, size, color, False, BODY)])
    text(s, x, y, w, h, runs, space=gap)


def stat(s, x, y, w, big, label, color=CLAY, big_size=34, sub=None):
    box(s, x, y, w, 1.35, fill=BG2, line=LINE)
    text(s, x+0.18, y+0.16, w-0.36, 0.7, [[R(big, big_size, color, True, MONO)]])
    text(s, x+0.18, y+0.86, w-0.36, 0.4, [[R(label.upper(), 10.5, INK2, True, MONO)]])
    if sub:
        text(s, x+0.18, y+1.06, w-0.36, 0.3, [[R(sub, 10, INK3, False, BODY)]])


def placeholder(s, x, y, w, h, label):
    shp = box(s, x, y, w, h, fill=BG2, line=CLAY, line_w=1.4)
    _set_dash(shp, "dash")
    text(s, x, y+h/2-0.5, w, 1.0,
         [[R("▣  ADD SCREENSHOT", 15, CLAY, True, MONO)],
          [R(label, 12, INK2, False, BODY)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def chip(s, x, y, w, txt, color=INK2, fill=PANEL, size=11, bold=True):
    box(s, x, y, w, 0.34, fill=fill, line=LINE)
    text(s, x, y+0.05, w, 0.26, [[R(txt, size, color, bold, MONO)]], align=PP_ALIGN.CENTER)


def connector(s, x1, y1, x2, y2, color=LINE, w=1.25, arrow=True):
    cn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    cn.line.color.rgb = color; cn.line.width = Pt(w)
    cn.shadow.inherit = False
    if arrow:
        lnEl = cn.line._get_or_add_ln()
        tail = lnEl.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'})
        lnEl.append(tail)
    return cn


# =========================================================================
# 1 — TITLE
s = slide(BG)
box(s, 0, 0, 13.333, 7.5, fill=BG, line=None)
text(s, 0.7, 2.2, 12, 0.5, [[R("ACCENTURE INNOVATION CHALLENGE 2026  ·  ROUND 2  ·  PROBLEM STATEMENT 4", 13, CLAY, True, MONO)]])
text(s, 0.68, 2.7, 12, 1.4, [[R("TwinForge", 66, INK, True, HEAD)]])
text(s, 0.7, 3.95, 12, 0.7, [[R("A live digital twin that ", 20, INK2, False, BODY),
                              R("prevents", 20, GREEN, True, BODY), R(", ", 20, INK2, False, BODY),
                              R("diagnoses", 20, BLUE, True, BODY), R(" and ", 20, INK2, False, BODY),
                              R("resolves", 20, CLAY, True, BODY),
                              R(" faults on a vehicle assembly line.", 20, INK2, False, BODY)]])
text(s, 0.7, 4.7, 12, 0.5, [[R("We don't mock a dashboard — we ", 15, INK2, False, BODY),
                             R("simulate a real factory with injectable cases", 15, INK, True, BODY),
                             R(", end to end.", 15, INK2, False, BODY)]])
box(s, 0.7, 5.7, 5.6, 0.9, fill=BG2, line=LINE)
text(s, 0.9, 5.82, 5.3, 0.7, [[R("Team TwinForge", 14, INK, True, BODY)],
                              [R("Sasanka Marthand · Sannith Reddy K · Pratyush Nayak B  —  IIT Madras", 11.5, INK2, False, BODY)]])

# =========================================================================
# 2 — THE PROBLEM
s = slide()
title(s, "The problem", "A vehicle line is a coupled flow system — and it fails two ways at once")
bullets(s, 0.6, 1.9, 6.3, 4.6, [
    ("Bottlenecks ripple. ", "A small capacity loss at one station grows queues, blocks upstream and starves downstream — the constraint moves ~6×/shift."),
    ("Defects surface late. ", "A drifting tool makes out-of-spec units with zero throughput change; they are only caught dozens of stations later at QA."),
    ("Sensor coverage is uneven. ", "Modern and legacy equipment mix — some stations are richly instrumented, others are manual / dark."),
    ("Dashboards are descriptive. ", "They show that a station is slow. They do not forecast where trouble forms, or trace a late symptom to its origin."),
    ("No live-control access. ", "A prototype cannot write set-points to safety-certified PLCs — it must advise, not act."),
], size=15, gap=13)
box(s, 7.3, 1.9, 5.4, 4.6, fill=BG2, line=LINE)
text(s, 7.55, 2.05, 5.0, 0.4, [[R("TWO INTELLIGENCE PROBLEMS", 12, CLAY, True, MONO)]])
box(s, 7.6, 2.55, 4.8, 1.7, fill=PANEL, line=GREEN)
text(s, 7.8, 2.68, 4.4, 1.5, [[R("PREVENTIVE", 13, GREEN, True, MONO)],
    [R("Forecast where a constraint will form from live telemetry — before it bites.", 13, INK, False, BODY)]])
box(s, 7.6, 4.4, 4.8, 1.85, fill=PANEL, line=BLUE)
text(s, 7.8, 4.53, 4.4, 1.6, [[R("DIAGNOSTIC", 13, BLUE, True, MONO)],
    [R("Trace a late symptom backward through the process graph to a ranked, probabilistic root cause.", 13, INK, False, BODY)]])

# =========================================================================
# 3 — THESIS
s = slide(BG2)
text(s, 0.6, 0.5, 12, 0.4, [[R("OUR THESIS", 13, CLAY, True, MONO)]])
text(s, 0.6, 1.0, 12.1, 2.0, [[R("A plant is instrumented for ", 30, INK, False, HEAD),
    R("utilisation", 30, AMBER, True, HEAD),
    R(" — but utilisation never tells you what is costing you cars.", 30, INK, False, HEAD)]])
text(s, 0.6, 3.05, 12, 0.9, [[R("TwinForge ranks stations by ", 17, INK2, False, BODY),
    R("what fixing them is worth, in cars", 17, INK, True, BODY),
    R(" — and it does so on a factory we actually simulate, case by case.", 17, INK2, False, BODY)]])
cards = [
    ("SIMULATE", "A discrete-event twin of a fixed 9-station network. Inject a case (fault + severity); the whole system responds live.", CLAY),
    ("PREVENT", "An ML forecast flags the station trending to constrain at T+5 — before failure. Test a fix by re-scoring the model.", GREEN),
    ("DIAGNOSE", "A defect with no flow signature is traced back with a probability on every station: genealogy × quality telemetry.", BLUE),
    ("RESOLVE", "When a bottleneck actually forms, re-simulate the fix under identical draws and report the cars recovered.", AMBER),
]
for i, (h, b, c) in enumerate(cards):
    x = 0.6 + i*3.06
    box(s, x, 4.2, 2.86, 2.6, fill=BG, line=c)
    text(s, x+0.2, 4.38, 2.5, 0.4, [[R(h, 14, c, True, MONO)]])
    text(s, x+0.2, 4.85, 2.5, 1.9, [[R(b, 12.5, INK2, False, BODY)]])

# =========================================================================
# 4 — THE FIXED FACTORY (DAG)
s = slide()
title(s, "The simulated factory", "A fixed 9-station DAG — deliberately non-linear, with heterogeneous sensors")
# draw the DAG
nodes = {
    "S1": (0.04,0.16,"Body Framing"), "S2": (0.04,0.50,"Chassis Prep"), "S3": (0.04,0.84,"Paint Shop"),
    "S4": (0.34,0.30,"Powertrain"), "S5": (0.34,0.58,"Interior"), "S6": (0.34,0.86,"Battery/HV"),
    "S7": (0.63,0.40,"Marriage"), "S8": (0.63,0.72,"Trim"), "S9": (0.90,0.55,"Final / QA"),
}
edges = [("S1","S4"),("S1","S5"),("S2","S4"),("S3","S5"),("S3","S6"),
         ("S4","S7"),("S6","S7"),("S5","S8"),("S7","S9"),("S8","S9")]
# region for the DAG
RX, RW = 0.9, 7.4
RY, RH = 2.0, 4.3
def px(nx): return RX + nx*RW
def py(ny): return RY + ny*RH
for (a,b) in edges:
    connector(s, px(nodes[a][0])+0.32, py(nodes[a][1])+0.22, px(nodes[b][0]), py(nodes[b][1])+0.22, color=LINE, w=1.2)
qual = {"S1","S3","S4","S6","S7","S8"}
for nid,(nx,ny,nm) in nodes.items():
    c = CLAY if nid=="S9" else (BLUE if nid in qual else LINE)
    b = box(s, px(nx), py(ny), 1.05, 0.44, fill=PANEL, line=c, line_w=1.4)
    text(s, px(nx), py(ny)+0.03, 1.05, 0.24, [[R(nid, 13, INK, True, MONO)]], align=PP_ALIGN.CENTER)
    text(s, px(nx), py(ny)+0.24, 1.05, 0.2, [[R(nm, 8.5, INK3, False, BODY)]], align=PP_ALIGN.CENTER)
text(s, 0.9, 6.55, 7.4, 0.3, [[R("blue = quality/defect origin  ·  clay = final QA sink  ·  edges = finite-capacity buffers", 10.5, INK3, False, MONO)]])
# right column facts
box(s, 8.7, 1.95, 4.05, 4.7, fill=BG2, line=LINE)
text(s, 8.9, 2.1, 3.7, 0.4, [[R("WHY THIS TOPOLOGY", 12, CLAY, True, MONO)]])
bullets(s, 8.9, 2.55, 3.7, 4.0, [
    ("3 sources, 1 sink. ", "5 balanced source→sink routes; merges & splits make propagation real."),
    ("Load varies. ", "S9 sees all flow, S7 three routes — so the constraint is not obvious."),
    ("Takt 12 s → 300 UPH. ", "Healthy line flows; a fault pushes one station over its own limit."),
    ("Heterogeneous sensors. ", "Torque, temp, voltage, vibration — some stations dark (manual)."),
], size=12.5, gap=9)

# =========================================================================
# 5 — ARCHITECTURE
s = slide()
title(s, "System architecture", "Pure-Python stack — simulation, ML and an honest arithmetic core")
steps = [
    ("DES SIMULATOR", "9-station DAG · CRN-seeded · conservation-exact", CLAY),
    ("FEATURES + ECON. TRUTH", "rolling windows · paired counterfactual labels", BLUE),
    ("ENGINES", "detector · forecast · diagnostic · counterfactual", GREEN),
    ("LIVE LOOP", "ingest→detect→forecast→emit · twin drift", AMBER),
    ("FASTAPI + SSE → UI", "vanilla JS + inline SVG · 4 workspaces", INK2),
]
y = 2.15; x = 0.7; bw = 2.28; gapx = 0.2
for i,(h,b,c) in enumerate(steps):
    xx = x + i*(bw+gapx)
    box(s, xx, y, bw, 1.5, fill=BG2, line=c)
    text(s, xx+0.14, y+0.16, bw-0.28, 0.5, [[R(h, 12.5, c, True, MONO)]])
    text(s, xx+0.14, y+0.62, bw-0.28, 0.8, [[R(b, 11, INK2, False, BODY)]])
    if i < len(steps)-1:
        connector(s, xx+bw, y+0.75, xx+bw+gapx, y+0.75, color=INK3, w=1.4)
box(s, 0.7, 4.1, 12.0, 2.5, fill=BG2, line=LINE)
text(s, 0.9, 4.25, 11.6, 0.4, [[R("DESIGN PRINCIPLE — DON'T TRAIN WHAT YOU CAN COMPUTE", 12, CLAY, True, MONO)]])
cols = [
    ("Arithmetic (no ML)", ["Blocked/starved boundary walk locates the current bottleneck", "Effective cycle time = processing ÷ availability", "Conservation, twin-drift, containment — all deterministic"], GREEN),
    ("Machine learning", ["Logistic forecast: constraint at T+5 from telemetry drift", "Defect backtrace: genealogy × quality-telemetry posterior", "Both trained against economic ground truth, not an identity"], BLUE),
]
for i,(h,items,c) in enumerate(cols):
    xx = 0.95 + i*5.95
    text(s, xx, 4.7, 5.6, 0.35, [[R(h, 14, c, True, BODY)]])
    bullets(s, xx, 5.1, 5.7, 1.4, items, size=12, gap=6, mark=c)

# =========================================================================
# 6 — FOUR WORKSPACES / FLOW
s = slide()
title(s, "The product", "Four IDE workspaces over one shared, injectable case")
flow = [("TELEMETRY","Watch the case run — flow, sensors, expected vs actual", CLAY),
        ("PREVENT","Model flags the risk early — test a fix (ML sensitivity)", GREEN),
        ("RESOLVE","Bottleneck formed? Re-simulate the fix, count the cars", AMBER)]
x=0.7; bw=3.7
for i,(h,b,c) in enumerate(flow):
    xx=x+i*(bw+0.35)
    box(s, xx, 2.1, bw, 1.5, fill=BG2, line=c)
    text(s, xx+0.18, 2.26, bw-0.36, 0.4, [[R(f"0{i+1}  {h}", 13, c, True, MONO)]])
    text(s, xx+0.18, 2.72, bw-0.36, 0.8, [[R(b, 12.5, INK2, False, BODY)]])
    if i<2: connector(s, xx+bw, 2.85, xx+bw+0.35, 2.85, color=INK3, w=1.6)
text(s, 0.7, 3.75, 12, 0.35, [[R("A single global CASE (fault + severity, injected from the top ribbon) drives all three.", 13, INK2, False, BODY)]])
box(s, 0.7, 4.35, 12.03, 2.25, fill=BG2, line=BLUE)
text(s, 0.9, 4.5, 11.6, 0.4, [[R("04  DIAGNOSE — DETACHED", 13, BLUE, True, MONO)]])
text(s, 0.9, 4.92, 11.6, 1.5, [[R("A defect has no throughput signature, so it stands apart. Pick a QA reject; the engine builds a probability graph — stations and connectors, no flow — assigning P(origin) to every station and lighting the ", 13.5, INK2, False, BODY),
    R("defect corridor", 13.5, BLUE, True, BODY),
    R(". A station on the corridor whose tool reads in spec is exonerated. Ships a containment list of affected VINs.", 13.5, INK2, False, BODY)]])

# =========================================================================
# 7 — TELEMETRY (screenshot)
s = slide()
title(s, "01 · Telemetry", "A Grafana-grade view of a line you can actually read")
bullets(s, 0.6, 1.95, 4.4, 4.4, [
    ("Live process flow. ", "Animated DAG — station state, WIP, buffer occupancy."),
    ("Expected vs actual. ", "Per-station throughput & cycle time against the design point."),
    ("Heterogeneous sensors. ", "Each station exposes only the sensors it has; manual stations are dark."),
    ("Drill-down. ", "Click any node for its full flow + sensor record with provenance."),
    ("Case-driven. ", "The top ribbon injects a fault; the dashboard responds in real time."),
], size=13.5, gap=11)
placeholder(s, 5.2, 1.95, 7.5, 4.5, "Telemetry tab — flow graph + expected/actual station table + forecast bars")

# =========================================================================
# 8 — PREVENT (screenshot)
s = slide()
title(s, "02 · Prevent", "Catch it before it bites — and prove the fix works, in milliseconds")
placeholder(s, 0.6, 1.95, 7.3, 4.5, "Prevent tab — pre-failure risk gauge + intervention what-if (before/after)")
bullets(s, 8.2, 1.95, 4.5, 4.4, [
    ("Pre-failure. ", "The forecast flags a station at T+5 while the line is still healthy."),
    ("ML sensitivity, not a sim. ", "Perturb telemetry (add operator / cool / throttle / service) → re-score the SAME model."),
    ("Measured. ", "e.g. S4 caught at ~min 9, risk 48% → cooling drops it to 0%."),
    ("Driving signals shown. ", "The exact telemetry the model reacts to (temp, current, vibration)."),
], size=13.5, gap=12)
stat(s, 8.2, 5.55, 2.15, "48→0%", "risk averted", GREEN, 24)
stat(s, 10.55, 5.55, 2.15, "T+5", "forecast horizon", CLAY, 24)

# =========================================================================
# 9 — DIAGNOSE (screenshot)
s = slide()
title(s, "03 · Diagnose", "A defect has no flow signature — so we build the graph backward")
placeholder(s, 0.6, 1.95, 7.3, 4.5, "Diagnose tab — probability graph (P per node) + defect corridor + ranked origins")
bullets(s, 8.2, 1.95, 4.5, 4.6, [
    ("The walk is blind here. ", "No throughput changes, so blocked/starved logic sees nothing."),
    ("Two independent signals. ", "Genealogy lift (which stations bad units share) × quality telemetry (which tool is out of spec)."),
    ("Disambiguation. ", "A bystander on the corridor with an in-spec tool is exonerated — S4 drops to 2% when S7 drifts."),
    ("Visual posterior. ", "P(origin) rendered on every one of 6 quality stations; corridor lit."),
    ("Containment. ", "The exact VINs built on the culprit since it began drifting."),
], size=13, gap=9)

# =========================================================================
# 10 — RESOLVE (screenshot)
s = slide()
title(s, "04 · Resolve", "A real counterfactual — simulation of a fix, over a simulation of the issue")
bullets(s, 0.6, 1.95, 4.4, 4.4, [
    ("Pure simulation. ", "The line is re-run under Common Random Numbers with the fix applied — no ML."),
    ("Paired & honest. ", "Current vs simulated on identical draws → cars gained is measured, not asserted."),
    ("Right vs wrong fix. ", "De-bottleneck the true constraint: +34 cars. Wrong station: +1."),
    ("Includes the free lever. ", "'Do nothing' and 'reduce release rate' (CONWIP) are options too."),
], size=13.5, gap=12)
placeholder(s, 5.2, 1.95, 7.5, 4.5, "Resolve tab — animated simulated line + current-vs-simulated compare + cumulative chart")

# =========================================================================
# 11 — RIGOR
s = slide(BG2)
title(s, "Engineering rigor", "We earned every number — and caught our own bugs")
bullets(s, 0.6, 1.95, 6.1, 4.6, [
    ("Economic ground truth. ", "Labels come from paired counterfactual re-simulation (which station's speed-up makes more cars), never the detector's own statistic — avoiding the identity trap."),
    ("Common Random Numbers. ", "Randomness is a deterministic function of (unit, station); pre-scheduled failures make counterfactuals truly paired."),
    ("Conservation-exact DES. ", "Injected = produced + WIP + backlog, asserted in tests."),
    ("Held-out & regret. ", "Forecast split by run; detector scored on cars lost, not top-1 accuracy."),
], size=13.5, gap=11)
box(s, 7.0, 1.95, 5.7, 4.6, fill=BG, line=LINE)
text(s, 7.25, 2.1, 5.2, 0.4, [[R("FOUR BUGS WE FOUND & FIXED", 12, CLAY, True, MONO)]])
bullets(s, 7.25, 2.55, 5.3, 4.0, [
    ("Resume-after-outage deadlock", " — a failed-mid-unit station froze forever."),
    ("Blocked-latch deadlock", " — one outage cascaded to a whole-line freeze; fix cut variance ±61→±1.6 cars."),
    ("CRN desync", " — shared RNG made a healthy speed-up 'gain' 68 cars; fix cut the noise floor to ~1."),
    ("Effective-CT double-counting", " — wall-clock proc time double-charged downtime; detector 0.62→1.00 top-1."),
], size=12.5, gap=9)

# =========================================================================
# 12 — RESULTS
s = slide()
title(s, "Measured results", "Reproducible bit-for-bit from fixed seeds")
row1 = [("~291", "UPH healthy (std 1.6)", CLAY),
        ("~1", "car noise floor (CRN)", GREEN),
        ("1.00", "detector top-1", BLUE),
        ("0.0", "detector regret (cars)", GREEN)]
row2 = [("92.5%", "forecast holdout (T+5)", CLAY),
        ("74%", "sparse-sensing agreement", AMBER),
        ("6", "defect origin stations", BLUE),
        ("0", "throughput cost of a defect", GREEN)]
for i,(b,l,c) in enumerate(row1):
    stat(s, 0.7+i*3.06, 2.1, 2.86, b, l, c, 34)
for i,(b,l,c) in enumerate(row2):
    stat(s, 0.7+i*3.06, 3.75, 2.86, b, l, c, 34)
box(s, 0.7, 5.55, 12.0, 1.15, fill=BG2, line=LINE)
text(s, 0.9, 5.7, 11.6, 0.9, [[R("Detector regret 0.0 cars vs 0.62 for naïve utilisation ranking  ·  every defect origin diagnosed correctly with the true culprit ranked #1  ·  full regeneration + retrain in one command.", 13, INK2, False, BODY)]])

# =========================================================================
# 13 — PS4 COVERAGE
s = slide()
title(s, "Coverage", "Every Problem-Statement-4 clause, answered by a working feature")
rows = [
    ("Bottleneck ripples downstream", "Blocked/starved walk + effective-CT detector"),
    ("Predict problems before they happen", "Logistic forecast of the T+5 constraint (Prevent)"),
    ("Defect uncaught for many vehicles", "Genealogy backtrace + containment list (Diagnose)"),
    ("Uneven / no sensor coverage", "Heterogeneous sensors + inference with confidence"),
    ("Multi-causal, intermittent faults", "Probabilistic posterior over origins, can be multiple"),
    ("No live-control modification", "The twin advises; it never writes set-points"),
    ("Validate vs outcomes / trust", "Economic ground truth, regret, twin-drift re-calibration"),
    ("Scale across lines & vintages", "Topology is data; engines are graph algorithms"),
]
y0 = 1.95; rh = 0.56
box(s, 0.7, y0, 12.0, 0.44, fill=PANEL, line=LINE)
text(s, 0.85, y0+0.08, 6.0, 0.3, [[R("PS4 CLAUSE", 11, CLAY, True, MONO)]])
text(s, 6.9, y0+0.08, 5.6, 0.3, [[R("TWINFORGE ANSWER", 11, CLAY, True, MONO)]])
for i,(a,b) in enumerate(rows):
    yy = y0+0.44+i*rh
    box(s, 0.7, yy, 12.0, rh, fill=BG2 if i%2 else BG, line=None)
    text(s, 0.85, yy+0.12, 6.0, 0.4, [[R(a, 12.5, INK, False, BODY)]])
    text(s, 6.9, yy+0.12, 5.7, 0.4, [[R(b, 12.5, INK2, False, BODY)]])

# =========================================================================
# 14 — TECH STACK + WHY WE WIN
s = slide(BG2)
title(s, "Why we win", "The only entry that is simultaneously simulated, rigorous and live")
comp = [
    ("Generic JS dashboards", "Polished, but no backend, no model, no simulation — the AI is vapour.", RED),
    ("Batch analytics scripts", "Rigorous numbers, but a digital shadow: no loop, no UI, no live case.", AMBER),
    ("TwinForge", "A running DES twin + trained ML + a live loop + four explainable workspaces.", GREEN),
]
for i,(h,b,c) in enumerate(comp):
    xx = 0.7 + i*4.06
    box(s, xx, 1.95, 3.85, 2.1, fill=BG, line=c)
    text(s, xx+0.18, 2.12, 3.5, 0.4, [[R(h, 13.5, c, True, BODY)]])
    text(s, xx+0.18, 2.62, 3.5, 1.4, [[R(b, 12.5, INK2, False, BODY)]])
text(s, 0.7, 4.35, 12, 0.4, [[R("TECH STACK", 12, CLAY, True, MONO)]])
bullets(s, 0.7, 4.8, 12.0, 1.8, [
    ("Backend. ", "Python 3.12 · custom discrete-event simulator · scikit-learn (logistic) · FastAPI + Server-Sent Events."),
    ("Frontend. ", "Zero-dependency vanilla JS + inline SVG (no Node build) — runs anywhere the Python backend does."),
    ("Reproducible. ", "Seeded generators, one-command data+model rebuild, a pytest invariant suite (10/10)."),
], size=13, gap=9)

# =========================================================================
# 15 — ROADMAP + CLOSE
s = slide(BG)
text(s, 0.6, 0.5, 12, 0.4, [[R("ROADMAP", 13, CLAY, True, MONO)]])
text(s, 0.6, 0.95, 12, 0.8, [[R("From a defensible prototype to a plant-ready twin", 30, INK, True, HEAD)]])
road = [
    ("Scale", "40+ stations across body / paint / final; segment-varying sensor coverage."),
    ("Retrofit optimiser", "Costed sensor-placement recommendation, phased by maintenance window."),
    ("Failure-mode classifier", "Wear vs sensor-drift vs material-lot, on scale-invariant ratios."),
    ("MLOps", "Drift monitoring, held-out recalibration, audit-grade replay of the decision surface."),
]
for i,(h,b) in enumerate(road):
    xx = 0.7 + (i%2)*6.1; yy = 2.1 + (i//2)*1.55
    box(s, xx, yy, 5.8, 1.35, fill=BG2, line=LINE)
    text(s, xx+0.2, yy+0.16, 5.4, 0.4, [[R(h, 15, CLAY, True, BODY)]])
    text(s, xx+0.2, yy+0.6, 5.4, 0.7, [[R(b, 12.5, INK2, False, BODY)]])
box(s, 0.7, 5.5, 12.0, 1.25, fill=BG2, line=CLAY)
text(s, 0.9, 5.62, 11.6, 1.0, [[R("TwinForge", 17, INK, True, HEAD),
    R("  —  we simulate the factory, forecast the failure, trace the defect, and price the fix. Live.", 15, INK2, False, BODY)]])
text(s, 0.9, 6.25, 11.6, 0.4, [[R("Thank you.", 14, CLAY, True, BODY)]])

# save
OUT = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "TwinForge_CaseDeck.pptx")
prs.save(OUT)
print("saved", OUT, "·", len(prs.slides._sldIdLst), "slides")
