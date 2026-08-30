"""
Build the TwinForge PITCH deck (python-pptx) - explanatory, tab-by-tab, with
both the high-level flow and the low-level detail (models, sensors, faults).
Run:  python -m scripts.build_pitch     (outputs TwinForge_Pitch.pptx)
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

BG=RGBColor(0x14,0x14,0x16); BG2=RGBColor(0x1C,0x1C,0x20); PANEL=RGBColor(0x22,0x22,0x28)
LINE=RGBColor(0x3A,0x3A,0x42); INK=RGBColor(0xED,0xED,0xED); INK2=RGBColor(0x9E,0x9E,0xA6)
INK3=RGBColor(0x6B,0x6B,0x74); CLAY=RGBColor(0xD9,0x77,0x57); GREEN=RGBColor(0x5F,0xB0,0x6A)
AMBER=RGBColor(0xD9,0xA4,0x41); RED=RGBColor(0xD0,0x68,0x5C); BLUE=RGBColor(0x7F,0xA8,0xD8)
HEAD="Calibri"; BODY="Calibri"; MONO="Consolas"

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
BLANK=prs.slide_layouts[6]

def slide(bg=BG):
    s=prs.slides.add_slide(BLANK); s.background.fill.solid(); s.background.fill.fore_color.rgb=bg; return s
def _dash(shape,d="dash"):
    ln=shape.line._get_or_add_ln(); e=ln.find(qn('a:prstDash'))
    if e is None: e=ln.makeelement(qn('a:prstDash'),{}); ln.append(e)
    e.set('val',d)
def box(s,x,y,w,h,fill=PANEL,line=LINE,lw=0.75,radius=True):
    shp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,Inches(x),Inches(y),Inches(w),Inches(h))
    if radius:
        try: shp.adjustments[0]=0.05
        except: pass
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb=fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb=line; shp.line.width=Pt(lw)
    shp.shadow.inherit=False; return shp
def R(t,sz=15,c=INK,b=False,f=BODY): return (t,sz,c,b,f)
def text(s,x,y,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,space=4,wrap=True):
    tb=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=tb.text_frame
    tf.word_wrap=wrap; tf.vertical_anchor=anchor
    tf.margin_left=0;tf.margin_right=0;tf.margin_top=0;tf.margin_bottom=0
    for i,para in enumerate(runs):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.alignment=align; p.space_after=Pt(space); p.space_before=Pt(0)
        for (t,sz,c,b,f) in para:
            r=p.add_run(); r.text=t; r.font.size=Pt(sz); r.font.bold=b; r.font.color.rgb=c; r.font.name=f
    return tb
def title(s,kicker,ttl):
    text(s,0.6,0.42,12.1,0.4,[[R(kicker.upper(),12,CLAY,True,MONO)]])
    text(s,0.6,0.74,12.1,0.9,[[R(ttl,30,INK,True,HEAD)]])
def bullets(s,x,y,w,h,items,size=14,gap=8,mark=CLAY,color=INK2):
    runs=[]
    for it in items:
        if isinstance(it,tuple):
            lead,rest=it; runs.append([R("▸  ",size,mark,True),R(lead,size,INK,True),R(rest,size,color)])
        else: runs.append([R("▸  ",size,mark,True),R(it,size,color)])
    text(s,x,y,w,h,runs,space=gap)
def stat(s,x,y,w,big,label,color=CLAY,bs=30):
    box(s,x,y,w,1.3,fill=BG2,line=LINE)
    text(s,x+0.16,y+0.15,w-0.32,0.7,[[R(big,bs,color,True,MONO)]])
    text(s,x+0.16,y+0.82,w-0.32,0.4,[[R(label.upper(),10,INK2,True,MONO)]])
def placeholder(s,x,y,w,h,label):
    shp=box(s,x,y,w,h,fill=BG2,line=CLAY,lw=1.4); _dash(shp)
    text(s,x,y+h/2-0.5,w,1.0,[[R("▣  ADD SCREENSHOT",15,CLAY,True,MONO)],[R(label,12,INK2)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
def conn(s,x1,y1,x2,y2,c=LINE,w=1.2,arrow=True):
    cn=s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,Inches(x1),Inches(y1),Inches(x2),Inches(y2))
    cn.line.color.rgb=c; cn.line.width=Pt(w); cn.shadow.inherit=False
    if arrow:
        ln=cn.line._get_or_add_ln(); ln.append(ln.makeelement(qn('a:tailEnd'),{'type':'triangle','w':'med','len':'med'}))
def table(s,x,y,w,rows,col_w,head_c=CLAY,rh=0.42):
    box(s,x,y,w,rh,fill=PANEL,line=LINE)
    cx=x
    for j,cell in enumerate(rows[0]):
        text(s,cx+0.12,y+0.09,col_w[j]-0.16,0.3,[[R(cell,10.5,head_c,True,MONO)]]); cx+=col_w[j]
    for i,row in enumerate(rows[1:]):
        yy=y+rh+i*rh; box(s,x,yy,w,rh,fill=BG2 if i%2 else BG,line=None); cx=x
        for j,cell in enumerate(row):
            bold = j==0
            text(s,cx+0.12,yy+0.1,col_w[j]-0.16,0.3,[[R(cell,11,INK if bold else INK2,bold)]]); cx+=col_w[j]

def tabslide(num,name,color,purpose,how,under,shot):
    s=slide(); title(s,f"Tab 0{num} · {name}",purpose)
    box(s,0.6,1.75,3.2,0.7,fill=BG2,line=color)
    text(s,0.78,1.9,2.9,0.4,[[R(f"WORKSPACE 0{num}",11,color,True,MONO)],[R(name,14,INK,True)]])
    text(s,0.6,2.7,6.0,0.35,[[R("HOW IT WORKS",11,CLAY,True,MONO)]])
    bullets(s,0.6,3.1,6.0,2.0,how,size=13,gap=8)
    text(s,0.6,5.15,6.0,0.35,[[R("UNDER THE HOOD",11,CLAY,True,MONO)]])
    bullets(s,0.6,5.5,6.0,1.6,under,size=12,gap=6,mark=color)
    placeholder(s,7.0,1.75,5.7,5.1,shot)

# 1 TITLE
s=slide(BG)
text(s,0.7,2.05,12,0.5,[[R("ACCENTURE INNOVATION CHALLENGE 2026 · ROUND 2 · PROBLEM STATEMENT 4",13,CLAY,True,MONO)]])
text(s,0.68,2.55,12,1.3,[[R("TwinForge",62,INK,True,HEAD)]])
text(s,0.7,3.75,12,0.6,[[R("A live digital twin of a vehicle assembly line — ",19,INK2),R("we simulate the factory",19,INK,True),R(", forecast the failure, trace the defect, and price the fix.",19,INK2)]])
for i,(h,c) in enumerate([("Telemetry",CLAY),("Prevent",GREEN),("Diagnose",BLUE),("Resolve",AMBER)]):
    box(s,0.7+i*3.05,4.7,2.85,0.65,fill=BG2,line=c); text(s,0.7+i*3.05,4.84,2.85,0.4,[[R(f"0{i+1}  {h}",14,c,True,MONO)]],align=PP_ALIGN.CENTER)
text(s,0.7,5.7,12,0.4,[[R("Team TwinForge · IIT Madras",13,INK2,True)]])

# 2 OVERVIEW
s=slide(); title(s,"Overview","One simulated line, one injectable case, four purpose-built workspaces")
bullets(s,0.6,1.85,6.3,3.0,[
    ("It is a simulator, not a mockup. ","A discrete-event twin of a fixed 9-station line runs live in the browser; every number is generated, not scripted."),
    ("A global Case drives everything. ","You inject a fault (station, type, severity, complication) once, from the ribbon — Telemetry, Prevent and Resolve all respond to it."),
    ("Four workspaces, four jobs. ","Watch (Telemetry), forecast & fix early (Prevent), trace a late defect (Diagnose), and re-simulate a fix (Resolve)."),
    ("Honest by design. ","Bottleneck location is arithmetic; ML is used only where it earns its place — forecasting and defect attribution."),
],size=14,gap=12)
flow=[("CASE","injected fault",CLAY),("TELEMETRY","live line",INK2),("PREVENT","forecast + fix",GREEN),("DIAGNOSE","defect trace",BLUE),("RESOLVE","re-simulate",AMBER)]
box(s,7.1,1.95,5.6,4.6,fill=BG2,line=LINE); text(s,7.3,2.1,5.2,0.4,[[R("THE LOOP",12,CLAY,True,MONO)]])
for i,(h,b,c) in enumerate(flow):
    yy=2.6+i*0.82; box(s,7.45,yy,4.9,0.62,fill=BG,line=c)
    text(s,7.6,yy+0.08,4.7,0.5,[[R(h+"  ",12.5,c,True,MONO),R(b,11.5,INK2)]])
    if i<len(flow)-1: conn(s,9.9,yy+0.62,9.9,yy+0.82,c=INK3,w=1.2)

# 3 THE FACTORY (DAG)
s=slide(); title(s,"What we simulate — the line","A fixed 9-station DAG with merges, splits and heterogeneous sensors")
nodes={"S1":(0.04,0.16,"Body Framing"),"S2":(0.04,0.50,"Chassis Prep"),"S3":(0.04,0.84,"Paint"),
 "S4":(0.34,0.30,"Powertrain"),"S5":(0.34,0.58,"Interior"),"S6":(0.34,0.86,"Battery"),
 "S7":(0.63,0.40,"Marriage"),"S8":(0.63,0.72,"Trim"),"S9":(0.90,0.55,"Final/QA")}
edges=[("S1","S4"),("S1","S5"),("S2","S4"),("S3","S5"),("S3","S6"),("S4","S7"),("S6","S7"),("S5","S8"),("S7","S9"),("S8","S9")]
RX,RW,RY,RH=0.9,7.2,1.95,4.2
px=lambda nx:RX+nx*RW; py=lambda ny:RY+ny*RH
qual={"S1","S3","S4","S6","S7","S8"}
for a,b in edges: conn(s,px(nodes[a][0])+0.32,py(nodes[a][1])+0.22,px(nodes[b][0]),py(nodes[b][1])+0.22,c=LINE,w=1.1)
for nid,(nx,ny,nm) in nodes.items():
    c=CLAY if nid=="S9" else (BLUE if nid in qual else LINE)
    box(s,px(nx),py(ny),1.02,0.44,fill=PANEL,line=c,lw=1.3)
    text(s,px(nx),py(ny)+0.03,1.02,0.24,[[R(nid,12.5,INK,True,MONO)]],align=PP_ALIGN.CENTER)
    text(s,px(nx),py(ny)+0.25,1.02,0.18,[[R(nm,8,INK3)]],align=PP_ALIGN.CENTER)
text(s,0.9,6.5,7.2,0.3,[[R("blue = quality/defect origin · clay = QA sink · edges = finite-capacity buffers",10,INK3,False,MONO)]])
box(s,8.6,1.95,4.1,4.7,fill=BG2,line=LINE); text(s,8.8,2.1,3.7,0.4,[[R("KEY FACTS",12,CLAY,True,MONO)]])
bullets(s,8.8,2.55,3.75,4.0,[
    ("3 sources → 1 sink. ","5 balanced routes; every unit sees 4 stations."),
    ("Takt 12 s → 300 UPH. ","Healthy line flows; a fault pushes one station over its limit."),
    ("Load varies. ","S9 sees all flow, S7 three routes — the constraint isn't obvious."),
    ("Heterogeneous sensors. ","Torque, temp, voltage, vibration; S2 fully manual (dark)."),
],size=12,gap=9)

# 4 SIMULATION DETAILS
s=slide(BG2); title(s,"What we simulate — the physics","A conservation-exact, seed-reproducible discrete-event model")
cols=[("Flow dynamics",[ "Per-second DES; units flow on routes","Blocked / starved / working / down states","Finite buffers → real back-pressure","Conservation: in = produced + WIP + backlog"],CLAY),
      ("Randomness (CRN)",[ "Noise is a function of (unit, station)","Failures pre-scheduled per station","So a counterfactual is truly paired","Noise floor ≈ 1 car"],BLUE),
      ("Faults injected",[ "Flow: gradual / sudden slowdown, outage","Complication: overheating, tool wear","Quality: tool-drift (defects, no flow hit)","Rare MTBF/MTTR background failures"],GREEN)]
for i,(h,items,c) in enumerate(cols):
    x=0.6+i*4.06; box(s,x,1.85,3.85,4.7,fill=BG,line=c)
    text(s,x+0.2,2.02,3.5,0.4,[[R(h,14,c,True)]])
    bullets(s,x+0.2,2.55,3.55,3.9,items,size=12.5,gap=11,mark=c)

# 5 WORKSPACE MAP
s=slide(); title(s,"The four workspaces","Each tab is its own view and toolset — like editors in an IDE")
cards=[("01 TELEMETRY","Watch the live line","Grafana-style dashboard: flow graph, per-station expected-vs-actual, heterogeneous sensors, drill-down.",CLAY),
       ("02 PREVENT","Forecast & fix early","ML flags a station trending to constrain at T+5 — test a fix by re-scoring the model. Pre-failure.",GREEN),
       ("03 DIAGNOSE","Trace a late defect","A QA reject with no flow signature; a probability graph attributes it to producer + handlers.",BLUE),
       ("04 RESOLVE","Re-simulate the fix","A real counterfactual: re-run the line with an intervention under identical draws; count cars.",AMBER)]
for i,(h,sub,b,c) in enumerate(cards):
    x=0.6+(i%2)*6.15; y=1.9+(i//2)*2.35
    box(s,x,y,5.9,2.15,fill=BG2,line=c)
    text(s,x+0.22,y+0.16,5.5,0.4,[[R(h,14,c,True,MONO)]])
    text(s,x+0.22,y+0.6,5.5,0.35,[[R(sub,14,INK,True)]])
    text(s,x+0.22,y+1.0,5.5,1.0,[[R(b,12.5,INK2)]])

# 6-9 TAB DETAILS
tabslide(1,"Telemetry",CLAY,"Purpose — read the living line at a glance, then drill in",
    [("Live flow graph. ","Animated DAG — state colour, WIP, buffer occupancy."),
     ("Expected vs actual. ","Per-station throughput & cycle time against the design point."),
     ("Heterogeneous sensors. ","Each station shows only the sensors it has; manual stations are dark."),
     ("Click to drill down. ","Full flow + sensor record with measured/inferred provenance.")],
    [("Source. ","Server-Sent-Events stream of the DES snapshot every few sim-seconds."),
     ("Forecast overlay. ","Amber ⚠ marks a station whose risk crosses 35% — before it constrains.")],
    "Telemetry tab — flow graph + station table (expected/actual + sensors) + forecast bars")
tabslide(2,"Prevent",GREEN,"Purpose — catch a forming constraint early and prove the fix",
    [("Pre-failure risk. ","The forecast flags a station at T+5 while the line is still healthy."),
     ("ML what-if (instant). ","Perturb telemetry (add operator / cool / throttle / service) → re-score the SAME model."),
     ("Coherent by cause. ","Cooling only helps overheating; servicing only helps wear; both gated to real sensors."),
     ("Driving signals shown. ","Exactly which telemetry the model reacts to.")],
    [("Model. ","Multinomial logistic regression over a 90-feature vector (9 stations × 10 channels)."),
     ("Not a simulation. ","The intervention is a feature perturbation re-scored in milliseconds.")],
    "Prevent tab — pre-failure risk gauge + intervention before/after (e.g. cool 48%→0%)")
tabslide(3,"Diagnose",BLUE,"Purpose — trace a late QA defect to its true origin",
    [("No flow signature. ","A defect changes no throughput, so the bottleneck walk is blind."),
     ("Attribute-aware. ","A coating reject could be S3 (applied) OR S7/S8 (scuffed it) — candidates = producer + handlers."),
     ("Distributed probability. ","Tempered posterior on every candidate node (e.g. 82 / 10 / 8%), not a false 98%."),
     ("Containment. ","The exact VINs built on the most-likely origin since it drifted.")],
    [("Signals fused. ","Quality telemetry (which tool is out of spec) × genealogy lift (which stations bad units share)."),
     ("Graph, not flow. ","Stations + connectors, no throughput — a pure probability view.")],
    "Diagnose tab — probability graph (P per node) + defect corridor + ranked producer/handler origins")
tabslide(4,"Resolve",AMBER,"Purpose — value a fix on a bottleneck that actually formed",
    [("Full re-simulation. ","The line is re-run under identical random draws with the fix applied."),
     ("Paired & measured. ","Current vs simulated on the same draws → cars gained is measured, not asserted."),
     ("Right vs wrong fix. ","De-bottleneck the true constraint: +34 cars; wrong station: +1."),
     ("Includes the free levers. ","'Do nothing' and 'reduce release rate' (CONWIP) are options.")],
    [("Engine. ","Common-Random-Numbers discrete-event re-run — no ML here."),
     ("Visibly running. ","The flow graph animates the fixed line as it re-simulates.")],
    "Resolve tab — animated simulated line + current-vs-simulated compare + cumulative-cars chart")

# 10 MODELS & METHODS
s=slide(); title(s,"Models & methods","Trained where it helps; computed where it doesn't")
rows=[["Engine","Method","Trained?"],
 ["Current constraint","Effective cycle time (proc ÷ availability) + blocked/starved walk","No — arithmetic"],
 ["Economic ground truth","Paired counterfactual re-sim (cars a speed-up actually yields)","No — measured"],
 ["Forecast (Prevent)","Multinomial logistic regression, 90 features, predict T+5 constraint","Yes"],
 ["Defect attribution","Genealogy lift × quality telemetry, tempered softmax over candidates","No — inference"],
 ["Counterfactual (Resolve)","CRN paired discrete-event re-simulation","No — simulation"]]
table(s,0.6,1.9,12.1,rows,[3.0,6.6,2.5])
text(s,0.6,5.6,12,1.0,[[R("The forecast is trained against ECONOMIC ground truth (which station's speed-up makes more cars), never the detector's own statistic — so a good score is a real forecast, not a self-fulfilling identity.",13,INK2)]])

# 11 SENSORS
s=slide(BG2); title(s,"Sensors simulated","Every station scans in/out throughput; internal coverage is heterogeneous")
rows=[["Station","Internal sensors","Quality channel"],
 ["S1 Body Framing","current · vibration · temp","weld"],
 ["S2 Chassis Prep","— (manual / dark)","—"],
 ["S3 Paint Shop","temp · current","coat"],
 ["S4 Powertrain","torque · current · temp · vibration","torque"],
 ["S5 Interior/Wiring","current · voltage","—"],
 ["S6 Battery / HV","voltage (near-dark)","cell"],
 ["S7 Underbody Marriage","torque · vibration · current","torque"],
 ["S8 Trim Line","torque · current","torque"],
 ["S9 Final / QA","temp · defect gate","—"]]
table(s,0.6,1.85,8.4,rows,[2.7,4.2,1.5],rh=0.46)
box(s,9.3,1.85,3.4,4.7,fill=BG,line=LINE); text(s,9.5,2.0,3.0,0.4,[[R("WHY IT MATTERS",11,CLAY,True,MONO)]])
bullets(s,9.5,2.5,3.05,4.0,[
    "A defect can only be confirmed where a sensor exists — dark stations force inference.",
    "Prevent only offers interventions a station's sensors support.",
    "Provenance (measured vs inferred) travels with every value.",
],size=12,gap=12)

# 12 FAULTS
s=slide(); title(s,"Faults & complications simulated","Each fault maps to a distinct signal — and a distinct fix")
rows=[["Fault","Moves","Fix that works"],
 ["Gradual / sudden slowdown","cycle time only","add operator (capacity)"],
 ["Outage (station down)","availability","planned stop / reroute"],
 ["Overheating (complication)","temperature ↑ (leading), then cycle","cool the machine"],
 ["Tool wear (complication)","vibration + current ↑, then cycle","service / re-seat tool"],
 ["Tool drift (quality)","quality channel out of spec (NO flow change)","re-calibrate the tool"]]
table(s,0.6,1.9,12.1,rows,[4.3,4.6,3.2],rh=0.5)
text(s,0.6,5.05,12,1.4,[[R("A mechanical slowdown does NOT heat a motor — so cooling it does nothing. Health signals rise only from a complication, and each is a ",13,INK2),R("leading indicator",13,INK,True),R(": the signal rises first, the slowdown follows, so Prevent can act before the constraint forms. This is what makes the fixes coherent instead of a black box that reacts to any change.",13,INK2)]])

# 13 RESULTS
s=slide(); title(s,"Measured results","Reproducible bit-for-bit from fixed seeds")
r1=[("~291","UPH healthy (std 1.6)",CLAY),("~1","car noise floor",GREEN),("1.00","detector top-1",BLUE),("0.0","detector regret (cars)",GREEN)]
r2=[("~89%","forecast holdout (T+5)",CLAY),("96%","sparse-sensing agreement",AMBER),("6","defect-origin stations",BLUE),("0","throughput cost of a defect",GREEN)]
for i,(b,l,c) in enumerate(r1): stat(s,0.7+i*3.06,2.0,2.86,b,l,c,32)
for i,(b,l,c) in enumerate(r2): stat(s,0.7+i*3.06,3.55,2.86,b,l,c,32)
box(s,0.7,5.3,12.0,1.2,fill=BG2,line=LINE)
text(s,0.9,5.45,11.6,0.9,[[R("128 simulated shifts · ~6,400 windowed samples split by run · every defect origin diagnosed with the true culprit ranked #1 · full data + model regenerate in one command.",13,INK2)]])

# 14 STACK / RUN
s=slide(BG2); title(s,"Tech stack & how to run","Pure Python + a zero-dependency web UI")
bullets(s,0.6,1.9,7.0,3.2,[
    ("Backend. ","Python 3.12 · custom discrete-event simulator · scikit-learn (logistic) · FastAPI + Server-Sent Events."),
    ("Frontend. ","Vanilla JS + inline SVG — no Node build; runs anywhere the Python backend does."),
    ("Reproducible. ","Seeded generators; one-command data+model rebuild; pytest invariant suite (10/10)."),
    ("Advises, never acts. ","The twin recommends; it never writes set-points to line control."),
],size=13.5,gap=12)
box(s,7.9,1.9,4.8,3.0,fill=BG,line=CLAY)
text(s,8.15,2.08,4.4,0.4,[[R("RUN IT",12,CLAY,True,MONO)]])
text(s,8.15,2.55,4.4,2.2,[[R("./run_demo.ps1",14,GREEN,True,MONO)],[R("then open",12,INK2)],[R("http://localhost:8000",13,INK,True,MONO)],
    [R("",6,INK2)],[R("regenerate model:",12,INK2)],[R("python -m scripts.\\",12,INK,False,MONO)],[R("generate_and_train --regen",12,INK,False,MONO)]])

# 15 CLOSE
s=slide(BG); text(s,0.6,2.4,12,0.9,[[R("Simulate. Forecast. Trace. Resolve.",34,INK,True,HEAD)]])
text(s,0.6,3.5,12,0.8,[[R("A digital twin that doesn't just visualise the line — it predicts the failure, finds the defect, and prices the fix, live.",17,INK2)]])
box(s,0.6,4.8,12.1,1.1,fill=BG2,line=CLAY)
text(s,0.85,5.0,11.6,0.8,[[R("TwinForge",17,INK,True,HEAD),R("  ·  Team TwinForge, IIT Madras  ·  Accenture Innovation Challenge 2026 — PS4",14,INK2)]])
text(s,0.85,6.4,11,0.4,[[R("Thank you.",14,CLAY,True)]])

OUT=os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))),"TwinForge_Pitch.pptx")
prs.save(OUT); print("saved",OUT,"·",len(prs.slides._sldIdLst),"slides")
